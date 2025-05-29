from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import JSONResponse
from typing import Optional, List
from uuid import uuid4
import os
import fitz  # PyMuPDF
import requests
from tempfile import NamedTemporaryFile
import cohere
from dotenv import load_dotenv
from sklearn.metrics.pairwise import cosine_similarity
from docx import Document
import openpyxl
from pptx import Presentation
from bs4 import BeautifulSoup

load_dotenv()
app = FastAPI()

# Cohere setup
COHERE_API_KEY = os.getenv("COHERE_API_KEY")
if not COHERE_API_KEY:
    raise ValueError("Missing COHERE_API_KEY in environment")
co = cohere.Client(COHERE_API_KEY)

session_store = {}

# ---------- Helpers ----------

def download_file(url):
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        ext = os.path.splitext(url)[-1]
        with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(response.content)
            return tmp.name
    except Exception as e:
        raise ValueError(f"Failed to download file: {e}")

def extract_text(path):
    ext = os.path.splitext(path)[-1].lower()
    try:
        if ext == ".pdf":
            doc = fitz.open(path)
            return "\n".join(page.get_text() for page in doc)

        elif ext == ".docx":
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)

        elif ext == ".txt":
            with open(path, "r", encoding="utf-8") as f:
                return f.read()

        elif ext in [".xlsx", ".xls"]:
            wb = openpyxl.load_workbook(path, data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text.append(" ".join(str(cell) for cell in row if cell is not None))
            return "\n".join(text)

        elif ext == ".pptx":
            prs = Presentation(path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        elif ext in [".html", ".htm"]:
            with open(path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f, "html.parser")
                return soup.get_text(separator="\n")

        else:
            raise ValueError(f"Unsupported file format: {ext}")

    except Exception as e:
        raise ValueError(f"Failed to read {ext} file: {e}")

def chunk_text(text, max_tokens=300):
    words = text.split()
    return [" ".join(words[i:i + max_tokens]) for i in range(0, len(words), max_tokens)]

def embed_chunks(chunks):
    response = co.embed(texts=chunks, model="embed-english-v3.0", input_type="search_document")
    return response.embeddings

def get_top_chunks(question, chunks, embeddings, top_k=5):
    query_embedding = co.embed(
        texts=[question], model="embed-english-v3.0", input_type="search_query"
    ).embeddings[0]
    sims = cosine_similarity([query_embedding], embeddings)[0]
    top_indices = sims.argsort()[::-1][:top_k]
    return [chunks[i] for i in top_indices]

def generate_answer(context, question):
    prompt = f"""You are a helpful assistant. Use the context below to answer the question strictly based on that context.

Context:
{context}

Question: {question}
Answer:"""
    response = co.generate(
        model="command-r-plus",
        prompt=prompt,
        max_tokens=500,
        temperature=0.3,
        stop_sequences=["--"]
    )
    return response.generations[0].text.strip()

# ---------- API Endpoints ----------

@app.post("/load")
async def load_multiple_files(
    file_urls: Optional[List[str]] = Form(None),
    files: Optional[List[UploadFile]] = File(None),
):
    try:
        if not file_urls and not files:
            return JSONResponse({"error": "Provide at least one file or URL."}, status_code=400)

        # Enforce max 8 total files/links
        total_inputs = (len(file_urls) if file_urls else 0) + (len(files) if files else 0)
        if total_inputs > 8:
            return JSONResponse({"error": "You can only upload a maximum of 8 files/URLs."}, status_code=400)

        all_text = []

        # Handle URLs
        if file_urls:
            for url in file_urls:
                downloaded = download_file(url)
                text = extract_text(downloaded)
                all_text.append(text)

        # Handle uploaded files
        if files:
            for file in files:
                ext = os.path.splitext(file.filename)[-1]
                with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                    tmp.write(await file.read())
                    file_path = tmp.name
                text = extract_text(file_path)
                all_text.append(text)

        combined_text = "\n".join(all_text)
        chunks = chunk_text(combined_text)
        embeddings = embed_chunks(chunks)

        session_id = str(uuid4())
        session_store[session_id] = {
            "chunks": chunks,
            "embeddings": embeddings
        }

        return JSONResponse({"message": "All documents processed.", "session_id": session_id})

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/chat")
async def chat_with_files(
    session_id: str = Form(...),
    question: str = Form(...),
):
    try:
        session = session_store.get(session_id)
        if not session:
            return JSONResponse({"error": "Invalid session_id."}, status_code=400)

        chunks = session["chunks"]
        embeddings = session["embeddings"]
        top_chunks = get_top_chunks(question, chunks, embeddings)
        context = "\n\n".join(top_chunks)
        answer = generate_answer(context, question)

        return JSONResponse({
            "session_id": session_id,
            "question": question,
            "answer": answer
        })

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

# ---------- Optional: Run ----------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000)
